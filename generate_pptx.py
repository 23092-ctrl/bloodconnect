from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

RED       = RGBColor(0xD3, 0x2F, 0x2F)
BLUE      = RGBColor(0x15, 0x65, 0xC0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x21, 0x21, 0x21)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEB)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
GREY      = RGBColor(0x75, 0x75, 0x75)
GREEN     = RGBColor(0x38, 0x8E, 0x3C)
ORANGE    = RGBColor(0xF5, 0x7C, 0x00)

IMG1 = '/home/cheikh/Downloads/demande/WhatsApp Image 2026-06-12 at 5.54.58 PM.jpeg'
IMG2 = '/home/cheikh/Downloads/demande/WhatsApp Image 2026-06-12 at 5.54.59 PM.jpeg'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # totalement vide

def slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_box(slide, items, x, y, w, h, size=14, color=DARK, icon='•'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f'{icon}  {item}'
        run.font.size = Pt(size)
        run.font.color.rgb = color

def red_bar(sl, h=1.0):
    rect(sl, 0, 0, 13.33, h, fill=RED)

def slide_number(sl, n, total=9):
    txt(sl, f'{n} / {total}', 12.5, 7.1, 0.8, 0.3, size=10, color=GREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RED)

# Goutte de sang (cercle blanc)
drop = sl.shapes.add_shape(9, Inches(5.9), Inches(0.5), Inches(1.5), Inches(1.5))
drop.fill.solid(); drop.fill.fore_color.rgb = WHITE
drop.line.fill.background()

txt(sl, 'BloodConnect', 1.5, 2.0, 10.33, 1.2,
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'Plateforme mobile intelligente de gestion des dons de sang',
    1.5, 3.3, 10.33, 0.7, size=20, color=RGBColor(0xFF,0xCC,0xCC),
    align=PP_ALIGN.CENTER)

rect(sl, 4.5, 4.3, 4.33, 0.06, fill=WHITE)

txt(sl, 'Juin 2026', 1.5, 4.6, 10.33, 0.5,
    size=14, color=RGBColor(0xFF,0xCC,0xCC), align=PP_ALIGN.CENTER)
txt(sl, 'github.com/23092-ctrl/bloodconnect', 1.5, 5.1, 10.33, 0.4,
    size=13, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Problématique', 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 2)

problems = [
    'Les centres de transfusion manquent de visibilité sur les donneurs disponibles',
    'Les donneurs ne savent pas où et quand donner, ni leur propre éligibilité',
    'Aucun système d\'alerte automatique lors de pénuries critiques',
    'Gestion manuelle des demandes de don, source d\'erreurs et de retards',
]
y = 1.4
for i, p in enumerate(problems):
    color = RED if i % 2 == 0 else BLUE
    rect(sl, 0.4, y, 0.06, 0.55, fill=color)
    txt(sl, p, 0.7, y, 12.3, 0.6, size=15, color=DARK)
    y += 0.8

rect(sl, 0.4, 5.0, 12.5, 1.0, fill=LIGHT_RED, line=RED)
txt(sl, '→  BloodConnect connecte donneurs et centres en temps réel avec alertes automatiques.',
    0.7, 5.1, 12.0, 0.7, size=15, bold=True, color=RED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — STACK TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Stack Technique', 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 3)

layers = [
    ('Mobile',         'Flutter 3.x',   'flutter_bloc · go_router · Dio 5',  BLUE),
    ('Backend',        'NestJS 10',      'TypeScript · Mongoose 8 · JWT',     RED),
    ('Base de données','MongoDB 7',      'Stockage documents NoSQL',           GREEN),
    ('Infrastructure', 'Docker Compose', 'Backend + MongoDB + Nginx',         ORANGE),
]

col_x = [0.4, 2.8, 5.5, 9.5]
rect(sl, 0.4, 1.3, 12.5, 0.45, fill=RGBColor(0xEE,0xEE,0xEE))
for label, cx in [('Couche',0.5),('Technologie',2.9),('Détails',5.6),('Rôle',9.6)]:
    txt(sl, label, cx, 1.33, 3.0, 0.38, size=12, bold=True, color=DARK)

y = 1.8
for layer, tech, detail, color in layers:
    bg = LIGHT_BG if y % 0.7 < 0.35 else WHITE
    rect(sl, 0.4, y, 12.5, 0.55, fill=bg)
    rect(sl, 0.4, y, 0.06, 0.55, fill=color)
    txt(sl, layer,  0.6,  y+0.05, 2.0, 0.45, size=13, bold=True,  color=color)
    txt(sl, tech,   2.9,  y+0.05, 2.4, 0.45, size=13, bold=True,  color=DARK)
    txt(sl, detail, 5.6,  y+0.05, 3.7, 0.45, size=12, color=GREY)
    y += 0.6

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Architecture', 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 4)

# Boites architecture
boxes = [
    (0.5,  1.4, 2.8, 4.5, BLUE,  'APPLICATION\nFLUTTER',
     ['Auth / Profile', 'Donations', 'Map & Alertes', 'flutter_bloc BLoC', 'go_router']),
    (4.2,  1.4, 2.8, 4.5, RED,   'API REST\nNESTJS',
     ['Auth JWT', 'Users', 'Appointments', 'Donations', 'Blood Centers', 'Notifications']),
    (7.9,  1.4, 2.0, 2.0, GREEN, 'MONGODB\n7',
     ['Documents', 'Mongoose ODM']),
    (7.9,  3.7, 2.0, 2.2, ORANGE,'NGINX\nProxy',
     ['Port 80/443', 'Reverse proxy', 'Docker']),
]
for bx, by, bw, bh, color, title, items in boxes:
    rect(sl, bx, by, bw, bh, line=color)
    rect(sl, bx, by, bw, 0.55, fill=color)
    txt(sl, title, bx+0.1, by+0.05, bw-0.2, 0.5,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, item in enumerate(items):
        txt(sl, f'• {item}', bx+0.15, by+0.65+i*0.38, bw-0.3, 0.35,
            size=10, color=DARK)

# Flèches
for label, x in [('HTTPS / JSON', 3.15), ('Docker\nnetwork', 7.25)]:
    txt(sl, '◄──────►', x, 3.2, 1.0, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)
    txt(sl, label, x-0.1, 3.55, 1.2, 0.4, size=9, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ROLES ET FONCTIONNALITES
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Rôles & Fonctionnalités', 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 5)

roles = [
    (RED,    'donor',         'Donneur',
     ['Soumettre une demande de don', 'Annuler une demande (pending/confirmed)',
      'Voir historique et statut', 'Alertes de pénurie sanguine']),
    (BLUE,   'center_admin',  'Admin Centre',
     ['Voir toutes les demandes du centre', 'Confirmer / Rejeter une demande',
      'Compléter un don (stock mis à jour)', 'Filtres par statut (onglets)']),
    (GREEN,  'admin',         'Admin Global',
     ['Accès tous centres', 'Gestion stocks sanguins',
      'Déclenchement alertes', 'Gestion utilisateurs']),
]

col_w = 4.1
for i, (color, role_id, label, items) in enumerate(roles):
    x = 0.4 + i * (col_w + 0.15)
    rect(sl, x, 1.2, col_w, 5.8, line=color)
    rect(sl, x, 1.2, col_w, 0.65, fill=color)
    txt(sl, label,   x+0.15, 1.25, col_w-0.3, 0.35, size=15, bold=True, color=WHITE)
    txt(sl, f'[{role_id}]', x+0.15, 1.55, col_w-0.3, 0.28, size=10, color=WHITE, italic=True)
    for j, item in enumerate(items):
        txt(sl, f'✓  {item}', x+0.2, 2.05+j*0.7, col_w-0.4, 0.6, size=12, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ELIGIBILITE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, "Règle d'Éligibilité Médicale", 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 6)

# Encadré principal
rect(sl, 0.5, 1.3, 12.33, 1.1, fill=LIGHT_RED, line=RED)
txt(sl, '🚫  Un donneur non éligible ne peut PAS soumettre de demande de don',
    0.8, 1.38, 11.8, 0.45, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(sl, 'jusqu\'à ce qu\'il redevienne éligible.',
    0.8, 1.78, 11.8, 0.4, size=14, color=RED, align=PP_ALIGN.CENTER)

# Deux cas
cases = [
    (RED,    '56 jours',        'Délai entre deux dons',
     ['Délai minimum de 56 jours (8 semaines) entre deux dons',
      'Calculé automatiquement depuis lastDonationDate',
      'L\'app affiche le nombre de jours restants',
      'Erreur 400 si la contrainte n\'est pas respectée']),
    (ORANGE, 'Non éligible',    'Statut médical',
     ['Champ medicallyEligible = false dans le profil',
      'Défini par un administrateur (contre-indication médicale)',
      'Bloque toute soumission jusqu\'à réactivation',
      'Visible sur la page Profil du donneur']),
]
for i, (color, badge, label, items) in enumerate(cases):
    x = 0.5 + i * 6.4
    rect(sl, x, 2.7, 6.2, 4.3, line=color)
    rect(sl, x, 2.7, 6.2, 0.55, fill=color)
    txt(sl, f'{badge}  —  {label}', x+0.15, 2.76, 5.9, 0.42,
        size=13, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, f'•  {item}', x+0.2, 3.4+j*0.62, 5.8, 0.55, size=12, color=DARK)

# Vérification double niveau
rect(sl, 0.5, 6.85, 12.33, 0.4, fill=RGBColor(0xE8,0xF5,0xE9))
txt(sl, '✓  Vérifié côté Frontend (bouton bloqué) ET côté Backend (erreur 400)',
    0.8, 6.88, 11.8, 0.3, size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CYCLE DE VIE D'UNE DEMANDE
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, "Cycle de Vie d'une Demande", 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 7)

states = [
    (ORANGE, 'PENDING',    1.0, 2.8),
    (GREEN,  'CONFIRMED',  4.5, 2.8),
    (BLUE,   'COMPLETED',  8.0, 2.8),
]
for color, label, x, y in states:
    rect(sl, x, y, 2.5, 0.9, fill=color)
    txt(sl, label, x, y+0.18, 2.5, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Flèches horizontales
for x in [3.6, 7.1]:
    txt(sl, '------>', x, 3.05, 0.9, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)

# REJECTED
rect(sl, 4.5, 4.5, 2.5, 0.9, fill=RED)
txt(sl, 'REJECTED', 4.5, 4.68, 2.5, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# CANCELLED
rect(sl, 1.0, 4.5, 2.5, 0.9, fill=GREY)
txt(sl, 'CANCELLED', 1.0, 4.68, 2.5, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Flèches vers le bas
for lx, color_arrow in [(2.1, GREY), (5.6, RED), (5.6, RED)]:
    txt(sl, '|', lx+0.6, 3.75, 0.5, 0.5, size=18, color=color_arrow, align=PP_ALIGN.CENTER)
    txt(sl, 'v', lx+0.6, 4.1, 0.5, 0.35, size=14, color=color_arrow, align=PP_ALIGN.CENTER)

# Légende
legend = [
    (ORANGE, 'Demande soumise par le donneur'),
    (GREEN,  'Confirmee par le centre'),
    (BLUE,   'Don effectue — stock mis a jour'),
    (RED,    'Refusee (avec motif optionnel)'),
    (GREY,   'Annulee par le donneur'),
]
for i, (color, label) in enumerate(legend):
    x = 0.5 + (i % 3) * 4.3
    y = 5.8 if i < 3 else 6.35
    rect(sl, x, y+0.1, 0.3, 0.3, fill=color)
    txt(sl, label, x+0.45, y+0.05, 3.7, 0.35, size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEMO / CAPTURES
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Démonstration — Interface Centre Admin', 0.5, 0.15, 12, 0.8,
    size=28, bold=True, color=WHITE)
slide_number(sl, 8)

# Insérer les deux captures
sl.shapes.add_picture(IMG1, Inches(0.4),  Inches(1.3), height=Inches(5.4))
sl.shapes.add_picture(IMG2, Inches(6.9),  Inches(1.3), height=Inches(5.4))

txt(sl, 'Fig. 1 — Demande PENDING\nBoutons Reject + Confirm',
    0.4, 6.75, 5.5, 0.6, size=11, color=GREY, align=PP_ALIGN.CENTER)
txt(sl, 'Fig. 2 — Demande CONFIRMED\nBoutons Reject + Complete',
    6.9, 6.75, 5.5, 0.6, size=11, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION / LIVRABLES
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
red_bar(sl, 1.1)
txt(sl, 'Livrables & Conclusion', 0.5, 0.15, 12, 0.8, size=32, bold=True, color=WHITE)
slide_number(sl, 9)

livrables = [
    (RED,   'Code source',     'github.com/23092-ctrl/bloodconnect'),
    (BLUE,  'APK Android',     'Build debug disponible — flutter build apk'),
    (GREEN, 'Rapport PDF',     'BloodConnect_Rapport_Technique.pdf'),
    (ORANGE,'Présentation',    'Ce fichier PowerPoint'),
]
for i, (color, label, detail) in enumerate(livrables):
    y = 1.4 + i * 0.85
    rect(sl, 0.4, y, 0.08, 0.6, fill=color)
    rect(sl, 0.5, y, 6.5, 0.6, fill=LIGHT_BG)
    txt(sl, label,  0.7, y+0.08, 2.5, 0.45, size=13, bold=True, color=color)
    txt(sl, detail, 3.3, y+0.08, 3.5, 0.45, size=12, color=GREY, italic=True)

rect(sl, 0.4, 5.0, 12.33, 1.5, fill=RED)
txt(sl, 'BloodConnect', 0.6, 5.1, 12.0, 0.55,
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'Relier chaque goutte de sang à ceux qui en ont besoin.',
    0.6, 5.6, 12.0, 0.4, size=14, color=RGBColor(0xFF,0xCC,0xCC),
    align=PP_ALIGN.CENTER, italic=True)
txt(sl, 'github.com/23092-ctrl/bloodconnect', 0.6, 6.0, 12.0, 0.35,
    size=12, color=WHITE, align=PP_ALIGN.CENTER)

out = '/home/cheikh/Downloads/BloodConnect_Presentation.pptx'
prs.save(out)
print('PPTX genere:', out)
